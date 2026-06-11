#!/usr/bin/env python3
"""Build the RAG Obsidian presentation deck per design.md (SNUBH Spine system).

Minimal & Light restyle (design.md §0): containers/cards/pills/boxes are outline
or no-fill — color lives in thin borders and text, never background fills; no top
accent bars; flat (no shadows); light Neutral Warm canvas everywhere. Charts keep
their data colors.

Narrative: pain → idea → SEE it (real note + workflow) → features → how (with
analogies) → why different → roadmap → try it → how-to ×2 → more tools → conclude
→ appendix Q&A ×2. 20 slides, auto page counter.
Run: python3 build_deck.py  → rag_obsidian_deck.pptx
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
os.makedirs(ASSETS, exist_ok=True)
LOGO = os.path.join(HERE, "logo.png")
VERSION = "v" + json.load(open(os.path.join(HERE, "..", "manifest.json")))["version"]
REPO = "github.com/grotyx/rag-obsidian"

# ---------------------------------------------------------------- palette
NEUTRAL_WARM = "f2f0eb"; WHITE = "ffffff"; HOUSE = "1E3932"; STARBUCKS = "006241"
ACCENT = "00754A"; GREEN_LIGHT = "d4e9e2"; UPLIFT = "2B5148"
GRAY = "8e8e93"; HAIR = "D9D9D9"; TEXT = "1F1F1F"; TEXT_SOFT = "6B6B6B"
NEUTRAL_COOL = "f9f9f9"; CERAMIC = "edebe9"; RED = "c82014"
FONT = "Pretendard"

def C(h): return RGBColor.from_string(h)

# ================================================================ charts
def build_charts():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.font_manager as fm
    import matplotlib.pyplot as plt
    base = "/Users/sangminpark/Library/CloudStorage/OneDrive-개인/08. Utilities/폰트/Pretendard-1.3.9/public/static/alternative"
    for w in ("Regular", "SemiBold", "Bold"):
        p = os.path.join(base, f"Pretendard-{w}.ttf")
        if os.path.exists(p):
            fm.fontManager.addfont(p)
    plt.rcParams["font.family"] = "Pretendard"
    plt.rcParams["axes.unicode_minus"] = False
    g_acc, gold = "#00754A", "#cba258"

    # ---- citation graph (그래프 slide) ----
    fig, ax = plt.subplots(figsize=(5.3, 4.5), dpi=200)
    have = {"2015\n딥러닝": (0.5, 0.78), "1998\nLeNet": (0.18, 0.5),
            "1997\nLSTM": (0.5, 0.30), "2017\nAlexNet": (0.84, 0.55)}
    edges = [("2015\n딥러닝", "1998\nLeNet"), ("2015\n딥러닝", "1997\nLSTM"),
             ("2017\nAlexNet", "1998\nLeNet")]
    for a, b in edges:
        (x1, y1), (x2, y2) = have[a], have[b]
        ax.plot([x1, x2], [y1, y2], color=g_acc, lw=2, zorder=1)
    for (x, y) in [(0.30, 0.88), (0.10, 0.80), (0.74, 0.86), (0.90, 0.30)]:
        ax.plot([0.5, x], [0.78, y], color=gold, lw=1.2, ls=":", zorder=1)
    for label, (x, y) in have.items():
        ax.scatter([x], [y], s=2400, c="#00754A", edgecolors="white", lw=2, zorder=2)
        ax.text(x, y, label, ha="center", va="center", fontsize=10, color="white", weight="bold", zorder=3)
    for label, (x, y) in [("놓친 ×3", (0.30, 0.88)), ("×2", (0.10, 0.80)), ("×2", (0.74, 0.86)), ("×2", (0.90, 0.30))]:
        ax.scatter([x], [y], s=1100, facecolors="none", edgecolors=gold, lw=1.8, ls=(0, (2, 1.5)), zorder=2)
        ax.text(x, y, label, ha="center", va="center", fontsize=8, color=gold, weight="bold", zorder=3)
    ax.text(0.5, 0.03, "● 보유 4편    ○ 놓친 논문 (9편 발굴 중 4편 표시)", ha="center", fontsize=10, color="#444")
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.axis("off")
    fig.tight_layout(pad=0.1); fig.savefig(f"{ASSETS}/graph.png", transparent=True); plt.close(fig)
    print("charts built →", ASSETS)

# ================================================================ pptx helpers
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

_PAGE = 0  # auto page counter — add_slide() bumps it, footer() reads it

def add_slide(bg=NEUTRAL_WARM):
    global _PAGE
    _PAGE += 1
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = C(bg); r.line.fill.background()
    r.shadow.inherit = False
    return s

def no_shadow(shape): shape.shadow.inherit = False

def set_run(run, size, bold=False, color=TEXT, italic=False, spc=True):
    run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = C(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {}); rPr.append(e)
        e.set("typeface", FONT)
    if spc:
        # -0.01em proportional: spc = -1 × font-pt centipoints (14pt → -14)
        rPr.set("spc", str(int(round(-1 * size))))

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf

def para(tf, runs, align=PP_ALIGN.LEFT, first=False, space_before=0, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before: p.space_before = Pt(space_before)
    if line: p.line_spacing = Pt(line)
    for spec in runs:
        r = p.add_run(); r.text = spec["t"]
        set_run(r, spec["s"], spec.get("b", False), spec.get("c", TEXT),
                spec.get("i", False))
    return p

def rounded(slide, x, y, w, h, fill=WHITE, line=HAIR, line_w=0.75, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try: sh.adjustments[0] = radius
    except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = C(line); sh.line.width = Pt(line_w)
    no_shadow(sh); return sh

def rect(slide, x, y, w, h, fill, line=None, line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = C(line); sh.line.width = Pt(line_w)
    no_shadow(sh); return sh

def num_badge(slide, cx, cy, n, d=0.44, fs=18):
    # Minimal & Light: outline circle, accent number — no fill
    ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy), Inches(d), Inches(d))
    ov.fill.background(); ov.line.color.rgb = C(ACCENT); ov.line.width = Pt(1.5)
    no_shadow(ov)
    _, tf = textbox(slide, cx - d / 2, cy, d, d, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [{"t": str(n), "s": fs, "b": True, "c": ACCENT}], align=PP_ALIGN.CENTER, first=True)

def step_arrow(s, x, y, w, h):
    _, tf = textbox(s, x, y, w, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [{"t": "→", "s": 24, "b": True, "c": ACCENT}], align=PP_ALIGN.CENTER, first=True)

def add_logo(slide):
    pic = slide.shapes.add_picture(LOGO, Inches(11.7545), Inches(0.0131), Inches(1.5788), Inches(0.8733))
    pic.crop_left = 0.0324; pic.crop_right = 0.1636; pic.crop_top = 0.2139; pic.crop_bottom = 0.1133
    return pic

# ---- zone builders
def header(slide, chapter):
    _, tf = textbox(slide, 0.1512, 0.135, 6.0, 0.28)
    para(tf, [{"t": chapter, "s": 11, "b": True, "c": GRAY}], first=True)

def headline(slide, text, color=STARBUCKS):
    _, tf = textbox(slide, 0.2752, 0.4771, 12.0, 0.7)
    para(tf, [{"t": text, "s": 30, "b": True, "c": color}], first=True, line=36)

def footer(slide, source=""):
    _, tf = textbox(slide, 0.1512, 7.24, 0.6, 0.25)
    para(tf, [{"t": str(_PAGE), "s": 9, "b": False, "c": GRAY}], first=True)
    if source:
        _, tf2 = textbox(slide, 7.1822, 7.2369, 6.0, 0.25)
        para(tf2, [{"t": source, "s": 8, "c": GRAY}], align=PP_ALIGN.RIGHT, first=True)

def card(slide, x=0.2752, y=1.4111, w=12.7946, h=5.7167):
    return rounded(slide, x, y, w, h, fill=WHITE, line=HAIR, line_w=0.75, radius=0.035)

def content(chapter, head, source=""):
    s = add_slide(NEUTRAL_WARM)
    header(s, chapter); headline(s, head); add_logo(s)
    card(s); footer(s, source)
    return s

# ================================================================ slide builders
def cover():
    s = add_slide(NEUTRAL_WARM); add_logo(s)
    rounded(s, 0.6, 1.11, 2.95, 0.49, fill=None, line=ACCENT, line_w=1.5, radius=0.5)
    _, tf = textbox(s, 0.6, 1.11, 2.95, 0.49, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [{"t": "OBSIDIAN PLUGIN", "s": 14, "b": True, "c": ACCENT}], align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(s, 0.6, 2.0512, 12.2953, 1.6)
    para(tf, [{"t": "Academic Paper Obsidian", "s": 40, "b": True, "c": STARBUCKS}], first=True, line=48)
    para(tf, [{"t": "Citation Manager", "s": 40, "b": True, "c": STARBUCKS}], line=48)
    rect(s, 0.6, 3.7593, 1.8, 0.05, fill=ACCENT)
    _, tf = textbox(s, 0.6, 4.0517, 11.5, 0.9)
    para(tf, [{"t": "내 논문 서재가 그대로 AI 검색엔진이 된다", "s": 20, "b": True, "c": UPLIFT}], first=True)
    para(tf, [{"t": "Obsidian 위에서 · 내 컴퓨터 안에서 · 답은 늘 출처와 함께", "s": 15, "c": TEXT_SOFT}], space_before=4)
    _, tf = textbox(s, 0.6, 4.98, 8.0, 0.35)
    para(tf, [{"t": REPO, "s": 14, "b": True, "c": ACCENT}], first=True)
    _, tf = textbox(s, 7.509, 5.58, 5.5, 0.32)
    para(tf, [{"t": "분당서울대학교병원 정형외과", "s": 16, "c": TEXT_SOFT}], align=PP_ALIGN.RIGHT, first=True)
    _, tf = textbox(s, 7.509, 5.95, 5.5, 0.5)
    para(tf, [{"t": "박상민", "s": 28, "b": True, "c": TEXT}], align=PP_ALIGN.RIGHT, first=True)
    _, tf = textbox(s, 0.2977, 6.9752, 8.0, 0.25)
    para(tf, [{"t": "연구실 세미나  |  2026.06", "s": 14, "c": TEXT_SOFT}], first=True)

def big_message(chapter, big_lines, subline):
    s = add_slide(NEUTRAL_WARM); header(s, chapter); add_logo(s); footer(s)
    _, tf = textbox(s, 0.6, 2.5, 12.1, 2.8, anchor=MSO_ANCHOR.MIDDLE)
    for i, ln in enumerate(big_lines):
        col = STARBUCKS if i == 0 else ACCENT
        para(tf, [{"t": ln, "s": 56, "b": True, "c": col}], align=PP_ALIGN.CENTER, first=(i == 0), line=64)
    _, tf = textbox(s, 0.9, 5.5, 11.5, 0.8)
    para(tf, [{"t": subline, "s": 16, "i": True, "c": TEXT_SOFT}], align=PP_ALIGN.CENTER, first=True, line=22)

def closing():
    s = add_slide(NEUTRAL_WARM); add_logo(s)
    _, tf = textbox(s, 0.6, 2.6, 12.1, 1.4, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [{"t": "Thank you", "s": 54, "b": True, "c": STARBUCKS}], align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(s, 0.6, 4.1, 12.1, 0.5)
    para(tf, [{"t": "Questions & discussion", "s": 18, "c": TEXT_SOFT}], align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(s, 0.6, 5.4, 12.1, 0.7)
    para(tf, [{"t": "박상민  ·  분당서울대학교병원 정형외과", "s": 14, "c": TEXT_SOFT}], align=PP_ALIGN.CENTER, first=True)
    para(tf, [{"t": REPO, "s": 13, "b": True, "c": ACCENT}], align=PP_ALIGN.CENTER, space_before=5)

# ---- patterns (inside body card)
CX, CY, CW, CH = 0.2752, 1.4111, 12.7946, 5.7167
IX, IY, IW = CX + 0.45, CY + 0.4, CW - 0.9

def takeaway_strip(s, text, y=None):
    y = y if y else CY + CH - 0.72
    rounded(s, CX + 0.35, y, CW - 0.7, 0.5, fill=None, line=ACCENT, line_w=1.5, radius=0.18)
    _, tf = textbox(s, CX + 0.6, y, CW - 1.2, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [{"t": text, "s": 14, "b": True, "c": STARBUCKS}], first=True)

def pattern_twoproblem(s, left, right, takeaway):
    gap = 0.5; cw = (IW - gap) / 2; cy = IY + 0.05; ch = 3.55
    for i, col in enumerate([left, right]):
        x = IX + i * (cw + gap)
        rounded(s, x, cy, cw, ch, fill=WHITE, line=HAIR, line_w=0.9, radius=0.05)
        _, tf = textbox(s, x + 0.4, cy + 0.32, 0.7, 0.7)
        para(tf, [{"t": "✗", "s": 32, "b": True, "c": RED}], first=True)
        _, tf = textbox(s, x + 0.4, cy + 1.08, cw - 0.8, 0.4)
        para(tf, [{"t": col["label"], "s": 14, "b": True, "c": TEXT_SOFT}], first=True)
        _, tf = textbox(s, x + 0.4, cy + 1.5, cw - 0.8, 1.1)
        para(tf, [{"t": col["problem"], "s": 24, "b": True, "c": TEXT}], first=True, line=30)
        _, tf = textbox(s, x + 0.4, cy + 2.72, cw - 0.8, 0.7)
        para(tf, [{"t": col["detail"], "s": 14, "c": TEXT_SOFT}], first=True, line=20)
    takeaway_strip(s, takeaway)

def note_mock(s, takeaway):
    px, py, pw, ph = IX, IY, 6.7, 4.3
    rounded(s, px, py, pw, ph, fill=WHITE, line=HAIR, line_w=0.9, radius=0.035)
    _, tf = textbox(s, px + 0.3, py + 0.22, pw - 0.6, 0.4)
    para(tf, [{"t": "References / ", "s": 13, "c": TEXT_SOFT},
              {"t": "lecun2015deep.md", "s": 13, "b": True, "c": TEXT}], first=True)
    rect(s, px + 0.3, py + 0.7, pw - 0.6, 0.014, fill=HAIR)
    _, tf = textbox(s, px + 0.35, py + 0.88, pw - 0.7, 2.1)
    fm = [("title", "Deep learning"), ("author", "LeCun Y · Bengio Y · Hinton G"),
          ("issued", "2015"), ("container-title", "Nature"),
          ("DOI", "10.1038/nature14539"), ("type", "article-journal")]
    first = True
    for k, v in fm:
        para(tf, [{"t": k + ":  ", "s": 13, "b": True, "c": ACCENT},
                  {"t": v, "s": 13, "c": TEXT}], first=first, line=20, space_before=0 if first else 2)
        first = False
    rect(s, px + 0.3, py + 3.18, pw - 0.6, 0.014, fill=HAIR)
    _, tf = textbox(s, px + 0.35, py + 3.32, pw - 0.7, 0.9)
    para(tf, [{"t": "## Abstract", "s": 12, "b": True, "c": TEXT_SOFT}], first=True, line=16)
    para(tf, [{"t": "Deep learning lets computational models learn representations of data with multiple levels of abstraction…",
               "s": 12, "c": TEXT_SOFT, "i": True}], space_before=2, line=16)
    ax = px + pw + 0.4; aw = IW - pw - 0.4
    anns = [("1", "파일 이름 = 인용키", "노트 한 개 = 논문 한 편"),
            ("2", "맨 위 = 메타데이터", "저자·연도·DOI를 표준 형식(CSL-JSON)으로. EndNote가 숨겨 두는 걸 사람이 읽는 평문으로."),
            ("3", "본문 = 검색·채팅이 읽는 글", "초록·메모·메인 텍스트")]
    ay = py + 0.08
    for nb, title, sub in anns:
        num_badge(s, ax + 0.22, ay, nb, d=0.42, fs=16)
        _, tf = textbox(s, ax + 0.62, ay - 0.06, aw - 0.62, 1.35)
        para(tf, [{"t": title, "s": 16, "b": True, "c": STARBUCKS}], first=True, line=20)
        para(tf, [{"t": sub, "s": 13, "c": TEXT_SOFT}], space_before=4, line=18)
        ay += 1.45
    takeaway_strip(s, takeaway)

def workflow_steps(s, steps, takeaway):
    n = len(steps); aw = 0.46
    cw = (IW - aw * (n - 1)) / n
    cy = IY + 0.18; ch = 3.4
    for i, st in enumerate(steps):
        x = IX + i * (cw + aw)
        rounded(s, x, cy, cw, ch, fill=WHITE, line=HAIR, line_w=0.9, radius=0.05)
        num_badge(s, x + cw / 2, cy + 0.26, i + 1)
        _, tf = textbox(s, x + 0.12, cy + 0.9, cw - 0.24, 0.6, anchor=MSO_ANCHOR.TOP)
        para(tf, [{"t": st["title"], "s": 15, "b": True, "c": STARBUCKS}], align=PP_ALIGN.CENTER, first=True, line=19)
        ebx = x + 0.18; eby = cy + 1.62; ebw = cw - 0.36; ebh = ch - 1.82
        rounded(s, ebx, eby, ebw, ebh, fill=None, line=HAIR, line_w=0.75, radius=0.08)
        _, tf = textbox(s, ebx + 0.1, eby + 0.1, ebw - 0.2, ebh - 0.2, anchor=MSO_ANCHOR.MIDDLE)
        first = True
        for ln in st["ex"]:
            para(tf, [{"t": ln["t"], "s": ln.get("s", 12), "b": ln.get("b", False),
                       "c": ln.get("c", TEXT), "i": ln.get("i", False)}],
                 align=PP_ALIGN.CENTER, first=first, line=16, space_before=0 if first else 4)
            first = False
        if i < n - 1:
            step_arrow(s, x + cw + 0.02, cy + ch / 2 - 0.3, aw - 0.04, 0.6)
    takeaway_strip(s, takeaway)

def pattern_C(s, claim, bullets, img, takeaway):
    lw = (CW - 0.9) * 0.52
    _, tf = textbox(s, IX, IY, lw, 3.6)
    para(tf, [{"t": claim, "s": 22, "b": True, "c": TEXT}], first=True, line=30)
    for b in bullets:
        p = tf.add_paragraph(); p.space_before = Pt(14); p.line_spacing = Pt(26)
        r = p.add_run(); r.text = "● "; set_run(r, 13, True, ACCENT)
        r2 = p.add_run(); r2.text = b; set_run(r2, 16, False, TEXT)
    rx = IX + lw + 0.3
    s.shapes.add_picture(img, Inches(rx), Inches(IY - 0.1), height=Inches(3.7))
    takeaway_strip(s, takeaway)

def pattern_E(s, cols, rows):
    nrow, ncol = len(rows) + 1, len(cols)
    tbl_w, tbl_h = CW - 0.9, 0.55 * nrow
    gtbl = s.shapes.add_table(nrow, ncol, Inches(IX), Inches(IY), Inches(tbl_w), Inches(tbl_h)).table
    gtbl.first_row = False; gtbl.horz_banding = False
    widths = [3.6] + [(tbl_w - 3.6) / (ncol - 1)] * (ncol - 1)
    for j, wv in enumerate(widths):
        gtbl.columns[j].width = Inches(wv)
    def fill_cell(cell, txt, size, bold, color, bg):
        cell.fill.solid(); cell.fill.fore_color.rgb = C(bg)
        cell.margin_left = Inches(0.1); cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt; set_run(r, size, bold, color)
    for j, ctxt in enumerate(cols):
        hi = (j == ncol - 1)
        fill_cell(gtbl.cell(0, j), ctxt, 12, True, STARBUCKS if hi else TEXT, NEUTRAL_COOL)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            hi = (j == ncol - 1)
            bg = WHITE if i % 2 == 0 else CERAMIC
            fill_cell(gtbl.cell(i + 1, j), val, 13,
                      bold=hi, color=(STARBUCKS if hi else (TEXT if j == 0 else TEXT_SOFT)), bg=bg)

def pattern_H(s, pillars, takeaway):
    n = len(pillars); gap = 0.22
    cw = (CW - 0.9 - gap * (n - 1)) / n
    cy, ch = IY + 0.1, 3.5
    for i, p in enumerate(pillars):
        x = IX + i * (cw + gap)
        rounded(s, x, cy, cw, ch, fill=WHITE, line=HAIR, line_w=0.75, radius=0.06)
        pw = min(cw - 0.4, 1.9)
        rounded(s, x + (cw - pw) / 2, cy + 0.28, pw, 0.34, fill=None, line=ACCENT, line_w=1.25, radius=0.5)
        _, tf = textbox(s, x + (cw - pw) / 2, cy + 0.28, pw, 0.34, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [{"t": p["tag"], "s": 11, "b": True, "c": ACCENT}], align=PP_ALIGN.CENTER, first=True)
        _, tf = textbox(s, x + 0.1, cy + 0.95, cw - 0.2, 1.7, anchor=MSO_ANCHOR.TOP)
        para(tf, [{"t": p["big"], "s": 34, "b": True, "c": STARBUCKS}], align=PP_ALIGN.CENTER, first=True)
        para(tf, [{"t": p["sub"], "s": 16, "b": True, "c": TEXT}], align=PP_ALIGN.CENTER, space_before=2)
        para(tf, [{"t": p["body"], "s": 12, "c": TEXT_SOFT}], align=PP_ALIGN.CENTER, space_before=10, line=16)
    takeaway_strip(s, takeaway)

def pattern_F(s, layers, takeaway):
    # Minimal & Light: outline bands, hierarchy via label text color
    label_cols = [(STARBUCKS, True), (TEXT, True), (ACCENT, True)]
    bh, gap = 1.05, 0.18
    by = IY
    for i, ly in enumerate(layers):
        lc, lb = label_cols[i % 3]
        rounded(s, IX, by, IW, bh, fill=None, line=HAIR, line_w=0.9, radius=0.05)
        _, tf = textbox(s, IX + 0.35, by, 2.0, bh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [{"t": ly["band"], "s": 15, "b": lb, "c": lc}], first=True)
        _, tf = textbox(s, IX + 2.4, by, IW - 2.8, bh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [{"t": ly["text"], "s": 15, "c": TEXT}], first=True, line=21)
        by += bh + gap
    takeaway_strip(s, takeaway, y=by + 0.02)

def pattern_J(s, milestones, takeaway):
    n = len(milestones)
    y_line = CY + 2.55
    x0, x1 = IX + 0.3, IX + IW - 0.3
    rect(s, x0, y_line, x1 - x0, 0.025, fill=ACCENT)
    step = (x1 - x0) / (n - 1)
    for i, m in enumerate(milestones):
        cx = x0 + i * step
        ov = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.11), Inches(y_line - 0.095), Inches(0.22), Inches(0.22))
        ov.fill.solid(); ov.fill.fore_color.rgb = C(ACCENT); ov.line.color.rgb = C(WHITE); ov.line.width = Pt(1.5)
        no_shadow(ov)
        _, tf = textbox(s, cx - 0.9, y_line - 0.9, 1.8, 0.5)
        para(tf, [{"t": m["year"], "s": 17, "b": True, "c": STARBUCKS}], align=PP_ALIGN.CENTER, first=True)
        _, tf = textbox(s, cx - 0.95, y_line + 0.25, 1.9, 1.0)
        para(tf, [{"t": m["label"], "s": 12, "c": TEXT}], align=PP_ALIGN.CENTER, first=True, line=15)
    takeaway_strip(s, takeaway)

def try_it(s, takeaway):
    steps = ["Obsidian 설치", "RAG Obsidian\n플러그인 활성화", "임베딩·LLM 키 설정"]
    n = 3; aw = 0.5; bw = (IW - aw * (n - 1)) / n; by = IY + 0.2; bh = 1.4
    for i, t in enumerate(steps):
        x = IX + i * (bw + aw)
        rounded(s, x, by, bw, bh, fill=None, line=ACCENT, line_w=1.5, radius=0.06)
        _, tf = textbox(s, x + 0.15, by, bw - 0.3, bh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [{"t": t, "s": 17, "b": True, "c": TEXT}], align=PP_ALIGN.CENTER, first=True, line=22)
        if i < n - 1:
            step_arrow(s, x + bw + 0.02, by + bh / 2 - 0.25, aw - 0.04, 0.5)
    uy = by + bh + 0.72
    rounded(s, IX + 0.8, uy, IW - 1.6, 1.05, fill=None, line=ACCENT, line_w=1.75, radius=0.1)
    _, tf = textbox(s, IX + 0.8, uy + 0.16, IW - 1.6, 0.35, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [{"t": "코드 · 문서 · 이슈", "s": 12, "b": True, "c": ACCENT}], align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(s, IX + 0.8, uy + 0.5, IW - 1.6, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [{"t": REPO, "s": 27, "b": True, "c": STARBUCKS}], align=PP_ALIGN.CENTER, first=True)
    _, tf = textbox(s, IX, uy + 1.25, IW, 0.4)
    para(tf, [{"t": f"현재 {VERSION} · 개발 빌드로 사용 중 · 공개 배포 예정", "s": 13, "c": TEXT_SOFT}], align=PP_ALIGN.CENTER, first=True)
    takeaway_strip(s, takeaway)

def usage_steps(s, rows, takeaway):
    ry = IY - 0.02; rh = 0.95; step = 1.08
    for i, r in enumerate(rows):
        y = ry + i * step
        rounded(s, IX, y, IW, rh, fill=WHITE, line=HAIR, line_w=0.9, radius=0.08)
        num_badge(s, IX + 0.42, y + (rh - 0.4) / 2, i + 1, d=0.4, fs=15)
        _, tf = textbox(s, IX + 0.75, y, 2.0, rh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [{"t": r["act"], "s": 15, "b": True, "c": STARBUCKS}], first=True, line=18)
        cx0 = IX + 2.85; cw0 = 4.6
        rounded(s, cx0, y + (rh - 0.55) / 2, cw0, 0.55, fill=None, line=HAIR, line_w=0.9, radius=0.12)
        _, tf = textbox(s, cx0 + 0.15, y + (rh - 0.55) / 2, cw0 - 0.3, 0.55, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [{"t": r["cmd"], "s": 12, "b": True, "c": ACCENT}], first=True, line=15)
        _, tf = textbox(s, cx0 + cw0 + 0.25, y, IX + IW - (cx0 + cw0) - 0.45, rh, anchor=MSO_ANCHOR.MIDDLE)
        for j, ln in enumerate(r["res"].split("\n")):
            para(tf, [{"t": ln, "s": 11.5, "c": TEXT_SOFT}], first=(j == 0), line=15,
                 space_before=0 if j == 0 else 2)
    takeaway_strip(s, takeaway)

def qa_rows(s, items, takeaway=None):
    n = len(items); gap = 0.22
    rh = (CH - 0.8 - gap * (n - 1)) / n
    y = IY
    for q, a in items:
        rounded(s, IX, y, IW, rh, fill=None, line=HAIR, line_w=0.9, radius=0.06)
        _, tf = textbox(s, IX + 0.4, y, IW - 0.8, rh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, [{"t": "Q.  ", "s": 15, "b": True, "c": STARBUCKS},
                  {"t": q, "s": 15, "b": True, "c": STARBUCKS}], first=True, line=20)
        para(tf, [{"t": "A.  ", "s": 12.5, "b": True, "c": ACCENT},
                  {"t": a, "s": 12.5, "c": TEXT}], space_before=6, line=17)
        y += rh + gap
    if takeaway:
        takeaway_strip(s, takeaway)

# ================================================================ assemble
def main():
    build_charts()
    cover()  # 1

    big_message("OPENING", ["답은 술술 나오는데", "출처가 가짜다"],
                "AI에 물으면 그럴듯한 답 — 인용은 지어낸 것. 정작 내 서재는 뜻으로 못 뒤진다.")  # 2

    s = content("BACKGROUND", "AI는 출처를 지어내고, 내 노트는 뜻으로 못 찾는다")  # 3
    pattern_twoproblem(s,
        {"label": "AI 챗봇에 물으면", "problem": "답은 그럴듯한데\n인용이 가짜다",
         "detail": "없는 DOI를 지어내고, 출처를 확인할 수 없다"},
        {"label": "내 노트·PDF를 뒤지면", "problem": "단어가 같아야\n겨우 찾힌다",
         "detail": "뜻이 같아도 표현이 다르면 놓친다 — 의미로 못 찾는다"},
        "필요한 건 — 내 서재 안에서, 뜻으로 찾아, 진짜 출처와 함께 답하는 도구")

    big_message("THE IDEA", ["내 논문 서재가", "곧 검색엔진이다"],
                "논문 1편 = 노트 1개. 그 노트들이 그대로 AI가 답하는 근거가 된다 — Zotero도, 서버도 없이.")  # 4

    s = content("HOW · 데이터", "노트 한 개가 곧 데이터베이스의 한 줄이다",
                "References/ 폴더의 마크다운 노트")  # 5
    note_mock(s, "이 평범한 .md 파일 하나 = 데이터베이스의 한 줄. 플러그인을 지워도 내 데이터는 남는다.")

    s = content("HOW · 흐름", "붙여넣기 한 번이면 출처 달린 답까지 간다",
                "실제 사용 흐름 예시")  # 6
    workflow_steps(s, [
        {"title": "DOI 붙여넣기", "ex": [{"t": "10.1038/nature14539", "s": 12.5, "b": True, "c": ACCENT}]},
        {"title": "노트 자동 생성", "ex": [{"t": "lecun2015deep.md", "s": 12.5, "b": True, "c": TEXT},
                                     {"t": "제목·저자·초록 자동 채움 ✓", "s": 11.5, "c": TEXT_SOFT}]},
        {"title": "서재에 질문", "ex": [{"t": "“표현 학습이", "s": 13, "b": True, "c": TEXT},
                                   {"t": "뭐였지?”", "s": 13, "b": True, "c": TEXT}]},
        {"title": "출처 달린 답", "ex": [{"t": "“…계층적 표현을", "s": 12, "c": TEXT},
                                    {"t": "학습한다 [1]”", "s": 12, "c": TEXT},
                                    {"t": "[1] LeCun 2015 → 노트", "s": 11.5, "b": True, "c": ACCENT}]},
    ], "붙여넣기 → 노트 → 질문 → 출처 달린 답. 서버 없이, 내 컴퓨터 안에서 한 흐름으로.")

    s = content("FEATURES", "한 플러그인 안에 네 가지: 서재·검색·채팅·그래프")  # 7
    pattern_H(s, [{"tag": "LIBRARY", "big": "서재", "sub": "Zotero 대체", "body": "DOI/PMID/arXiv\n→ 마크다운 노트"},
                  {"tag": "SEARCH", "big": "검색", "sub": "단어+뜻", "body": "하이브리드\n+ 메타데이터 필터"},
                  {"tag": "CHAT", "big": "채팅", "sub": "근거 기반", "body": "구절 인용 [n]\n→ 서식 출처"},
                  {"tag": "GRAPH", "big": "그래프", "sub": "OpenAlex", "body": "관련 ·\n놓친 논문"}],
              "추가로: PDF 임포트 · @인용 자동완성 · 참고문헌 자동생성 · 온톨로지(분류체계) 팩")

    s = content("HOW · 검색", "검색은 단어와 뜻을 함께 본다",
                "하이브리드 검색 · 임베딩 provider 3종")  # 8
    pattern_F(s, [{"band": "단어로", "text": "정확한 용어 매칭 — 똑똑해진 Ctrl+F (BM25)"},
                  {"band": "뜻으로", "text": "의미가 가까우면 매칭 — 'MI'로 검색해도 'myocardial infarction'을 찾음 (벡터 임베딩)"},
                  {"band": "+ 필터", "text": "저자·연도·저널로 좁히기 (노트 맨 위 메타데이터)"}],
              "단어 + 뜻 + 메타데이터를 한 번에 — 전부 로컬·무료. 필요하면 인용 그래프로 확장.")

    s = content("HOW · 채팅", "모든 답은 내 서재의 진짜 구절을 인용한다",
                "근거 기반(RAG) 채팅 파이프라인")  # 9
    pattern_F(s, [{"band": "① 검색", "text": "질문 → 내 서재에서 관련 구절만 회수 (서재 밖은 보지 않음)"},
                  {"band": "② 근거", "text": "LLM이 그 구절로만 답하고, 문장마다 [n] 표시 — 출처 없는 주장 금지"},
                  {"band": "③ 인용", "text": "[n] → APA·Vancouver 서식으로, 클릭하면 원문 노트로 이동"}],
              "슬라이드 2의 '가짜 인용' 문제 해결 — 모든 [n]은 내 서재의 진짜 구절")

    s = content("HOW · 그래프", "인용 그래프가 '내가 놓친 논문'을 찾아준다",
                "테스트 vault 4편 · OpenAlex (2.5억 편) 라이브")  # 10
    pattern_C(s, "OpenAlex 인용 데이터로\n서재의 '인용 지도'를 그린다",
              ["내 논문들이 서로 인용하는 관계 = 지도의 선",
               "여러 번 인용되는데 내겐 없는 논문 = 놓친 논문",
               "다음에 뭘 읽을지 — LLM 비용 0으로"],
              f"{ASSETS}/graph.png", "실데이터: 4편만 넣어도 핵심 누락 논문 9편 자동 발굴")

    s = content("WHY DIFFERENT", "서지 도구도 AI 도구도 못 하던 일을 한 번에",
                "Smart Connections 5.1k★ · Copilot 7.1k★ · Citation 1.3k★")  # 11
    pattern_E(s, ["기능", "서지 플러그인", "AI/RAG 플러그인", "RAG Obsidian"],
              [["의미(뜻) 검색", "✗", "✓", "✓"], ["저자·연도 필터", "✓", "✗", "✓"],
               ["구절 단위 인용", "✗", "✗ (노트 단위)", "✓"], ["서식 참고문헌", "일부", "✗", "✓"],
               ["인용 그래프", "✗", "✗", "✓"], ["Zotero 불필요", "✗", "—", "✓"]])

    s = content("ROADMAP", "Phase 0부터 5까지, 각 단계가 따로도 쓸모 있다",
                f"{VERSION} · {REPO}")  # 12
    pattern_J(s, [{"year": "P0", "label": "서지관리\nDOI→노트"}, {"year": "P1", "label": "의미 검색\nOrama"},
                  {"year": "P2", "label": "근거 채팅"}, {"year": "P3", "label": "PDF\n임포트"},
                  {"year": "P4", "label": "인용\n그래프"}, {"year": "P5", "label": "작성 지원\n참고문헌"}],
              "향후: 풀 CSL(citeproc) 스타일 · 온톨로지 검색확장 · 모바일 QA")

    s = content("GET STARTED", "Obsidian만 있으면 오늘부터 시작")  # 13
    try_it(s, "Obsidian만 있으면 오늘부터 — 내 서재가 답하기 시작한다")

    s = content("HOW TO USE ①", "서재 채우는 법 — 네 가지 입구",
                "명령 팔레트 = Cmd/Ctrl+P · 'RAG'만 쳐도 전부 검색")  # 14
    usage_steps(s, [
        {"act": "DOI·PMID 추가", "cmd": "“Add reference by DOI / PMID / arXiv”",
         "res": "식별자 붙여넣기 → References/에 노트 자동 생성.\n제목·저자·연도·초록까지 자동으로 채워진다"},
        {"act": "PubMed 검색", "cmd": "“Search PubMed and add references”  (리본: 돋보기)",
         "res": "키워드 검색 → 체크해서 한 번에 추가.\n옵션: LLM 자동 요약(OA는 전문 기반) + MeSH 태그"},
        {"act": "PDF 던져넣기", "cmd": "“Import PDF into library”",
         "res": "PDF에서 텍스트·DOI를 찾아 메타데이터 자동 인식.\n식별자 없으면 LLM이 제목·저자를 추출"},
        {"act": "Zotero·EndNote 이관", "cmd": "“Import references (BibTeX / RIS / CSL-JSON)”",
         "res": "기존 서재를 내보내기 파일로 한 번에 가져온다.\n중복은 자동으로 건너뜀"},
    ], "넣는 방법이 무엇이든 결과는 같다 — 노트 하나. 처음 한 번만 “Rebuild search index”.")

    s = content("HOW TO USE ②", "찾고, 묻고, 쓰는 법 — 매일 쓰는 네 가지",
                "명령 팔레트 = Cmd/Ctrl+P · 'RAG'만 쳐도 전부 검색")  # 15
    usage_steps(s, [
        {"act": "뜻으로 검색", "cmd": "“Search library (semantic)”",
         "res": "단어+뜻 하이브리드 — 'MI'로 'myocardial infarction'도 잡힌다.\n저자·연도·저널 필터 병행"},
        {"act": "서재에 질문", "cmd": "“Chat with library”  (리본: 말풍선)",
         "res": "내 서재 구절만 근거로 답하고 문장마다 [n].\n[n] 클릭 → 원문 노트로 이동"},
        {"act": "글 쓰며 인용", "cmd": "본문에서  @  입력",
         "res": "제목·저자로 검색해 [@citekey] 삽입.\n읽기 화면에선 (Park, 2022)·[1] 같은 서식 인용으로 렌더"},
        {"act": "참고문헌 생성", "cmd": "“Update bibliography in current note”",
         "res": "## References를 저널 스타일로 자동 생성.\n노트에 csl: 한 줄로 저널 교체 — 1만+ 스타일 자동 다운로드"},
    ], "추가 → 검색 → 질문 → 인용 → 참고문헌, 전부 Cmd+P 안에서.")

    s = content("MORE TOOLS", "추가부터 투고까지, 연구 워크플로 전체를 덮는다",
                "명령 30개 · 전부 Cmd+P")  # 16
    pattern_H(s, [{"tag": "EXPORT", "big": "내보내기", "sub": "lock-in 0", "body": "BibTeX/RIS/CSL-JSON로\n언제든 통째로 도로\n가져갈 수 있다"},
                  {"tag": "SUMMARIZE", "big": "요약", "sub": "LLM + MeSH", "body": "추가 시 자동 요약·MeSH\n(OA는 전문 기반)"},
                  {"tag": "WRITE", "big": "투고", "sub": "원고 컴파일", "body": "[@키] 전부 풀어\n인용+참고문헌 완성본\n(Pandoc-ready)"},
                  {"tag": "CURATE", "big": "관리", "sub": "서재 위생", "body": "철회 경고 · 중복 정리\n읽기 큐 · OA PDF 다운로드\nPDF 하이라이트 추출"}],
              "lock-in 0 (언제든 내보내기) — 들어온 뒤엔 읽기·쓰기·투고까지 한 앱에서")

    s = content("CONCLUSION", "Zotero 없이, 서버 없이, 늘 출처와 함께")  # 17
    pattern_H(s, [{"tag": "OWN YOUR DATA", "big": "마크다운", "sub": "= 내 DB", "body": "평문·git·로컬 우선\n플러그인보다 오래 남음"},
                  {"tag": "GROUNDED", "big": "근거", "sub": "구절 인용", "body": "출처 없는 답 없음\n서식 참고문헌까지"},
                  {"tag": "GRAPH-LITE", "big": "그래프", "sub": "놓친 논문", "body": "OpenAlex 인용엣지\nLLM 비용 0"}],
              "내 서재가 곧 검색엔진이고, 답은 늘 출처와 함께 — RAG Obsidian")

    s = content("APPENDIX · Q&A", "자주 나올 질문 ①")  # 18
    qa_rows(s, [
        ("Zotero랑 뭐가 다른가요?",
         "Zotero는 별도 앱+별도 DB. 여기선 글 쓰는 Obsidian 노트 자체가 서지 DB라 한 곳에서 끝나고, AI 근거 검색이 기본 내장."),
        ("오프라인에서도 되나요?",
         "검색·색인은 전부 로컬. 임베딩도 로컬(Ollama) 가능. 네트워크는 메타데이터 가져오기와 인용 그래프(OpenAlex)뿐."),
        ("내 데이터는 어디에 저장되나요?",
         "전부 내 vault 안 평문 마크다운. 외부 서버 없음 — 내가 키를 넣은 LLM 호출만 예외."),
    ])

    s = content("APPENDIX · Q&A", "자주 나올 질문 ②")  # 19
    qa_rows(s, [
        ("SNOMED·MeSH 같은 의학 용어체계도 되나요?",
         "온톨로지 팩(JSON)으로 장착. IS_A 계층·동의어 링크 지원. PubMed 추가 시 MeSH 태그는 기본."),
        ("수천 편짜리 큰 서재도 버티나요?",
         "수백~수천 편 타깃. 색인이 인메모리라 아주 큰 vault는 추후 sqlite-vec 교체 예정."),
        ("가짜 인용(환각)은 정말 없나요?",
         "답을 서재 구절에만 근거하게 강제하고 [n]은 실제 노트에 연결. 근거 없으면 \"근거 없음\"이라 답하게 설계."),
    ])

    closing()  # 20

    out = os.path.join(HERE, "rag_obsidian_deck.pptx")
    prs.save(out)
    print("saved →", out, "·", len(prs.slides), "slides")

if __name__ == "__main__":
    main()
